"""
Document Agent —— 解析 PDF / Word / PPT / TXT / Markdown 文档，
并将提取文本送入 LLM → 知识图谱。
"""
import os
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import markdown as md_lib
from src.agents.base import BaseAgent
from src.services.llm_service import LLMService
from src.services.kg_service import KGService
from src.services.rag_service import RAGService
from typing import Dict, Any, Optional
import logging

logger = logging.getLogger(__name__)


class DocumentAgent(BaseAgent):
    """多格式文档解析 + 知识提取 Agent"""

    AGENT_TYPE = "document"

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}

    def __init__(self, llm=None, kg=None, rag=None):
        super().__init__(llm=llm, kg=kg, rag=rag)

    # ------- 公共入口 -------
    async def process(self, file_path: str) -> Optional[Dict[str, Any]]:
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            logger.warning(f"不支持的文件格式: {ext}")
            return None

        parser = getattr(self, f"_parse_{ext[1:]}", None)
        if not parser:
            return None

        raw_text = parser(file_path)
        if not raw_text:
            logger.warning(f"文件内容为空: {file_path}")
            return None

        self._ensure_services()

        # LLM 提取
        extraction = await self.llm.extract_knowledge(raw_text[:10000])

        # 入库
        source_info = {"type": ext, "path": file_path}
        self.kg.upsert_knowledge(extraction, source_info)
        self.rag.add_documents([raw_text], [{"source": file_path, "type": ext}])

        logger.info(f"✅ 文档解析完成: {file_path}")
        return extraction

    # ------- 各格式解析器 -------
    @staticmethod
    def _parse_pdf(path: str) -> str:
        doc = fitz.open(path)
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text.strip()

    @staticmethod
    def _parse_docx(path: str) -> str:
        doc = DocxDocument(path)
        return "\n".join(p.text for p in doc.paragraphs if p.text).strip()

    @staticmethod
    def _parse_pptx(path: str) -> str:
        prs = Presentation(path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text.strip():
                            lines.append(p.text.strip())
        return "\n".join(lines)

    @staticmethod
    def _parse_txt(path: str) -> str:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()

    @staticmethod
    def _parse_md(path: str) -> str:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()


# 自动注册
from src.agents.base import agent_registry
agent_registry.register(DocumentAgent)
